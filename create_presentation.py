#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Кольори для дизайну
    BLUE = RGBColor(41, 128, 185)
    DARK_BLUE = RGBColor(52, 73, 94)
    LIGHT_GRAY = RGBColor(236, 240, 241)
    DARK_GRAY = RGBColor(52, 73, 94)
    WHITE = RGBColor(255, 255, 255)
    GREEN = RGBColor(39, 174, 96)
    ORANGE = RGBColor(230, 126, 34)

    # === СЛАЙД 1: Титульний ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "📸 Photo Gallery"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Веб-застосунок для пошуку та перегляду фотографій"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = LIGHT_GRAY
    subtitle_para.alignment = PP_ALIGN.CENTER

    author_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1.5))
    author_frame = author_box.text_frame
    author_text = "Виконав: Брам Олександр\nГрупа: ОП-221\nОлександрійський політехнічний фаховий коледж\n2026"
    author_frame.text = author_text
    author_para = author_frame.paragraphs[0]
    author_para.font.size = Pt(16)
    author_para.font.color.rgb = LIGHT_GRAY
    author_para.alignment = PP_ALIGN.CENTER

    # === СЛАЙД 2: Головна сторінка сайту ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    tf.text = "Головна сторінка застосунку"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Placeholder для скріншота
    img_placeholder = slide.shapes.add_shape(
        1, Inches(0.5), Inches(1.2), Inches(9), Inches(5)
    )
    img_placeholder.fill.solid()
    img_placeholder.fill.fore_color.rgb = LIGHT_GRAY
    img_placeholder.line.color.rgb = DARK_GRAY
    img_placeholder.line.width = Pt(2)

    img_text = img_placeholder.text_frame
    img_text.text = "📷 СКРІНШОТ:\n\nГоловна сторінка з навігацією,\nформою пошуку та галереєю фотографій"
    for para in img_text.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(20)
        para.font.color.rgb = DARK_GRAY

    # Опис знизу
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.4), Inches(9), Inches(0.9))
    tf = desc_box.text_frame
    tf.text = "🔹 Navbar з назвою застосунку  🔹 Форма пошуку з кнопкою  🔹 Адаптивна галерея фотографій (3 колонки)"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

    # === СЛАЙД 3: Форма пошуку ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    tf.text = "🔍 Пошук фотографій"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Ліворуч - скріншот
    img_placeholder = slide.shapes.add_shape(
        1, Inches(0.5), Inches(1.2), Inches(5.5), Inches(4)
    )
    img_placeholder.fill.solid()
    img_placeholder.fill.fore_color.rgb = LIGHT_GRAY
    img_placeholder.line.color.rgb = DARK_GRAY
    img_placeholder.line.width = Pt(2)

    img_text = img_placeholder.text_frame
    img_text.text = "📷 СКРІНШОТ:\n\nФорма пошуку з\nвведеним запитом\n\"nature\""
    for para in img_text.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(18)
        para.font.color.rgb = DARK_GRAY

    # Праворуч - опис
    desc_box = slide.shapes.add_textbox(Inches(6.2), Inches(1.2), Inches(3.5), Inches(4.5))
    tf = desc_box.text_frame
    text = """Функціонал пошуку:

✅ Текстове поле з
   placeholder

✅ HTML5 валідація
   (required)

✅ JavaScript валідація
   (trim для порожніх)

✅ Кнопка "Пошук"

✅ Submit при Enter"""

    tf.text = text
    for para in tf.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = DARK_GRAY
        para.space_after = Pt(8)

    # === СЛАЙД 4: Категорії ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    tf.text = "🗂 Категорії для швидкого доступу"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Скріншот
    img_placeholder = slide.shapes.add_shape(
        1, Inches(0.5), Inches(1.2), Inches(9), Inches(2.5)
    )
    img_placeholder.fill.solid()
    img_placeholder.fill.fore_color.rgb = LIGHT_GRAY
    img_placeholder.line.color.rgb = DARK_GRAY
    img_placeholder.line.width = Pt(2)

    img_text = img_placeholder.text_frame
    img_text.text = "📷 СКРІНШОТ: Горизонтальна навігація з кнопками категорій\n(Featured, Wallpapers, 3D Renders, Nature, Textures, Film, Architecture...)"
    for para in img_text.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(18)
        para.font.color.rgb = DARK_GRAY

    # Опис
    desc_box = slide.shapes.add_textbox(Inches(0.8), Inches(4), Inches(8.5), Inches(3))
    tf = desc_box.text_frame
    text = """🎯 11 популярних категорій для швидкого пошуку

📱 Горизонтальна прокрутка (на мобільних - свайп)

🎨 Активна категорія виділена темним фоном та білим текстом

⚡ Клік на категорію автоматично виконує пошук"""

    tf.text = text
    for para in tf.paragraphs:
        para.font.size = Pt(20)
        para.font.color.rgb = DARK_GRAY
        para.space_after = Pt(12)

    # === СЛАЙД 5: Галерея фотографій ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    tf.text = "🖼 Адаптивна галерея фотографій"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Скріншот
    img_placeholder = slide.shapes.add_shape(
        1, Inches(0.5), Inches(1.2), Inches(9), Inches(4.5)
    )
    img_placeholder.fill.solid()
    img_placeholder.fill.fore_color.rgb = LIGHT_GRAY
    img_placeholder.line.color.rgb = DARK_GRAY
    img_placeholder.line.width = Pt(2)

    img_text = img_placeholder.text_frame
    img_text.text = "📷 СКРІНШОТ:\n\nГалерея з 9-12 фотографіями в адаптивній сітці (3 колонки)\nКожна картка має кнопку улюбленого (🤍), інформацію про автора та лайки"
    for para in img_text.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(18)
        para.font.color.rgb = DARK_GRAY

    # Опис
    desc_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.9), Inches(8.5), Inches(1.4))
    tf = desc_box.text_frame
    text = """Bootstrap Grid: row-cols-1 row-cols-md-3 g-4  •  Динамічне створення карток через createElement  •  Lazy loading зображень"""
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

    # === СЛАЙД 6: Фільтри та сортування ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    tf.text = "🔄 Сортування та фільтри"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Ліворуч - скріншот
    img_placeholder = slide.shapes.add_shape(
        1, Inches(0.5), Inches(1.2), Inches(5), Inches(3.5)
    )
    img_placeholder.fill.solid()
    img_placeholder.fill.fore_color.rgb = LIGHT_GRAY
    img_placeholder.line.color.rgb = DARK_GRAY
    img_placeholder.line.width = Pt(2)

    img_text = img_placeholder.text_frame
    img_text.text = "📷 СКРІНШОТ:\n\nФорми фільтрів:\n- Сортування\n- Мінімум лайків"
    for para in img_text.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(16)
        para.font.color.rgb = DARK_GRAY

    # Праворуч - опис
    desc_box = slide.shapes.add_textbox(Inches(5.7), Inches(1.2), Inches(4), Inches(5))
    tf = desc_box.text_frame
    text = """📊 Сортування:
• Релевантність
• Від найпопулярніших ⬇️
• Від найновіших 🕒

🎯 Фільтр лайків:
• Мінімум вподобань
• Number input (min="0")
• Debouncing 500мс
• Тільки для улюблених

⚡ Автоматичне
   оновлення галереї"""

    tf.text = text
    for para in tf.paragraphs:
        para.font.size = Pt(15)
        para.font.color.rgb = DARK_GRAY
        para.space_after = Pt(6)

    # === СЛАЙД 7: Режими завантаження ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    tf.text = "📄 Три режими завантаження"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Скріншот
    img_placeholder = slide.shapes.add_shape(
        1, Inches(0.5), Inches(1.2), Inches(9), Inches(2)
    )
    img_placeholder.fill.solid()
    img_placeholder.fill.fore_color.rgb = LIGHT_GRAY
    img_placeholder.line.color.rgb = DARK_GRAY
    img_placeholder.line.width = Pt(2)

    img_text = img_placeholder.text_frame
    img_text.text = "📷 СКРІНШОТ: Перемикач режимів (btn-group)\n📄 Пагінація  |  ➕ Завантажити більше  |  ∞ Infinite Scroll"
    for para in img_text.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(16)
        para.font.color.rgb = DARK_GRAY

    # Опис режимів
    desc_box = slide.shapes.add_textbox(Inches(0.7), Inches(3.5), Inches(8.7), Inches(3.5))
    tf = desc_box.text_frame
    text = """📄 Пагінація - класична навігація з номерами сторінок
   • Кнопки "Попередня"/"Наступна"
   • Максимум 5 видимих номерів
   • Плавна прокрутка до верху

➕ Завантажити більше - кнопка для додавання фото
   • Додає наступні 12 фото до існуючих
   • Контроль завантаження користувачем

∞ Infinite Scroll - автоматичне завантаження при прокрутці
   • За 300px до кінця сторінки
   • Debouncing 200мс
   • Безперервний перегляд"""

    tf.text = text
    for para in tf.paragraphs:
        para.font.size = Pt(15)
        para.font.color.rgb = DARK_GRAY
        para.space_after = Pt(8)

    # === СЛАЙД 8: Модальне вікно ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    tf.text = "🖼 Перегляд фото у повному розмірі"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Скріншот
    img_placeholder = slide.shapes.add_shape(
        1, Inches(0.5), Inches(1.2), Inches(9), Inches(4.5)
    )
    img_placeholder.fill.solid()
    img_placeholder.fill.fore_color.rgb = LIGHT_GRAY
    img_placeholder.line.color.rgb = DARK_GRAY
    img_placeholder.line.width = Pt(2)

    img_text = img_placeholder.text_frame
    img_text.text = "📷 СКРІНШОТ:\n\nМодальне вікно Bootstrap з фото у високій якості\nЗаголовок, опис, автор, кількість лайків\nКнопки: \"Закрити\" та \"Відкрити оригінал\""
    for para in img_text.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(18)
        para.font.color.rgb = DARK_GRAY

    # Опис
    desc_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.9), Inches(8.5), Inches(1.4))
    tf = desc_box.text_frame
    text = """Bootstrap Modal (modal-xl)  •  Клік на картку відкриває вікно  •  Закриття: ESC, клік поза вікном, кнопка ✕"""
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

    # === СЛАЙД 9: Улюблені фото ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    tf.text = "❤️ Система улюблених фотографій"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Скріншот
    img_placeholder = slide.shapes.add_shape(
        1, Inches(0.5), Inches(1.2), Inches(9), Inches(4)
    )
    img_placeholder.fill.solid()
    img_placeholder.fill.fore_color.rgb = LIGHT_GRAY
    img_placeholder.line.color.rgb = DARK_GRAY
    img_placeholder.line.width = Pt(2)

    img_text = img_placeholder.text_frame
    img_text.text = "📷 СКРІНШОТ:\n\nВкладка \"❤️ Улюблені (N)\" зі збереженими фотографіями\nКнопки улюбленого: 🤍 (не додано) та ❤️ (додано)\nФільтр за мінімальною кількістю лайків"
    for para in img_text.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(18)
        para.font.color.rgb = DARK_GRAY

    # Опис
    desc_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.4), Inches(8.5), Inches(1.8))
    tf = desc_box.text_frame
    text = """💾 localStorage для збереження даних  •  JSON.stringify / JSON.parse  •  Лічильник улюблених\nКлік на ❤️ додає/видаляє  •  Фільтрація за лайками  •  Дані зберігаються навіть після закриття браузера"""
    p = tf.paragraphs[0]
    p.font.size = Pt(15)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

    # === СЛАЙД 10: Адаптивний дизайн ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    tf.text = "📱 Адаптивний дизайн"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Три placeholder для різних екранів
    # Desktop
    desktop = slide.shapes.add_shape(
        1, Inches(0.5), Inches(1.2), Inches(3), Inches(4)
    )
    desktop.fill.solid()
    desktop.fill.fore_color.rgb = LIGHT_GRAY
    desktop.line.color.rgb = DARK_GRAY
    desktop.line.width = Pt(2)

    dt = desktop.text_frame
    dt.text = "📷\n\nDesktop\n\n3 колонки"
    for para in dt.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(16)
        para.font.color.rgb = DARK_GRAY

    # Tablet
    tablet = slide.shapes.add_shape(
        1, Inches(3.7), Inches(1.2), Inches(3), Inches(4)
    )
    tablet.fill.solid()
    tablet.fill.fore_color.rgb = LIGHT_GRAY
    tablet.line.color.rgb = DARK_GRAY
    tablet.line.width = Pt(2)

    tt = tablet.text_frame
    tt.text = "📷\n\nTablet\n\n2 колонки"
    for para in tt.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(16)
        para.font.color.rgb = DARK_GRAY

    # Mobile
    mobile = slide.shapes.add_shape(
        1, Inches(6.9), Inches(1.2), Inches(3), Inches(4)
    )
    mobile.fill.solid()
    mobile.fill.fore_color.rgb = LIGHT_GRAY
    mobile.line.color.rgb = DARK_GRAY
    mobile.line.width = Pt(2)

    mt = mobile.text_frame
    mt.text = "📷\n\nMobile\n\n1 колонка"
    for para in mt.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(16)
        para.font.color.rgb = DARK_GRAY

    # Опис
    desc_box = slide.shapes.add_textbox(Inches(0.7), Inches(5.4), Inches(8.7), Inches(1.8))
    tf = desc_box.text_frame
    text = """Bootstrap Grid System: row-cols-1 row-cols-md-3\nMedia queries для різних пристроїв\nГнучкі елементи (Flexbox) для категорій та фільтрів\nHover ефекти на картках (тільки desktop)"""
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

    # === СЛАЙД 11: Обробка помилок ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    tf.text = "⚠️ Обробка помилок"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Скріншот
    img_placeholder = slide.shapes.add_shape(
        1, Inches(0.5), Inches(1.2), Inches(5.5), Inches(3)
    )
    img_placeholder.fill.solid()
    img_placeholder.fill.fore_color.rgb = RGBColor(255, 230, 230)
    img_placeholder.line.color.rgb = RGBColor(200, 50, 50)
    img_placeholder.line.width = Pt(2)

    img_text = img_placeholder.text_frame
    img_text.text = "📷 СКРІНШОТ:\n\nBootstrap Alert\nз повідомленням\nпро помилку"
    for para in img_text.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(16)
        para.font.color.rgb = RGBColor(180, 40, 40)

    # Опис
    desc_box = slide.shapes.add_textbox(Inches(6.2), Inches(1.2), Inches(3.5), Inches(5.5))
    tf = desc_box.text_frame
    text = """Типи помилок:

🔴 API помилки:
• 401 - Невірний ключ
• 403 - Ліміт запитів

🔴 Мережа:
• Помилка з'єднання
• Timeout

🔴 Валідація:
• Порожній запит

🔴 localStorage:
• Недоступний

Try-catch блоки
для всіх операцій"""

    tf.text = text
    for para in tf.paragraphs:
        para.font.size = Pt(14)
        para.font.color.rgb = DARK_GRAY
        para.space_after = Pt(6)

    # === СЛАЙД 12: Технічний стек ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    tf.text = "🛠 Технічний стек"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Ліва колонка
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(4.3), Inches(5.8))
    tf = left_box.text_frame
    text = """⚡ Збірка:
• Vite 7.2.4
• npm/Node.js

💻 Frontend:
• JavaScript ES6+
• HTML5
• CSS3

📚 Бібліотеки:
• Bootstrap 5.3.8
• Axios 1.13.2"""

    tf.text = text
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = DARK_GRAY
        paragraph.space_after = Pt(10)

    # Права колонка
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(5.8))
    tf = right_box.text_frame
    text = """🔌 API:
• Unsplash API
• REST

💾 Дані:
• localStorage
• JSON

🛠 Інструменти:
• Git/GitHub
• VS Code
• Chrome DevTools"""

    tf.text = text
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = DARK_GRAY
        paragraph.space_after = Pt(10)

    # === СЛАЙД 13: Модульна архітектура ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    tf.text = "📁 Модульна архітектура"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.5), Inches(5.8))
    tf = content_box.text_frame
    text = """Проєкт організовано за принципом розділення відповідальностей:

📁 src/config.js - конфігурація API ключів та констант
📁 src/state.js - глобальний стан застосунку
📁 src/dom.js - посилання на всі DOM елементи
📁 src/storage.js - робота з localStorage (getFavorites, saveFavorites)
📁 src/api.js - HTTP запити до Unsplash API через Axios
📁 src/ui.js - функції відображення UI (createPhotoCard, displayPhotos)
📁 src/events.js - обробники подій (submit, click, scroll)
📁 src/main.js - точка входу, ініціалізація застосунку

✅ Переваги: легко підтримувати, розширювати, тестувати
✅ Кожен модуль має свою відповідальність
✅ Можливість повторного використання коду"""

    tf.text = text
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = DARK_GRAY
        paragraph.space_after = Pt(8)

    # === СЛАЙД 14: Приклад коду ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    tf.text = "💻 Приклад коду: API запити"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    code_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.5), Inches(5.8))
    tf = code_box.text_frame
    code = """export async function fetchPhotos(query, page) {
  try {
    const response = await axios.get(
      `${UNSPLASH_API_URL}/search/photos`,
      {
        params: {
          query: query,
          page: page,
          per_page: 12,
          client_id: UNSPLASH_ACCESS_KEY
        }
      }
    )

    return response.data.results
  } catch (error) {
    if (error.response) {
      showError(`Помилка API: ${error.response.status}`)
    } else if (error.request) {
      showError('Помилка мережі')
    }
    return []
  }
}"""

    tf.text = code
    for paragraph in tf.paragraphs:
        paragraph.font.name = 'Courier New'
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = DARK_GRAY

    # === СЛАЙД 15: Результати ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    tf.text = "🎯 Результати та досягнення"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.5), Inches(5.8))
    tf = content_box.text_frame
    text = """✅ Повністю реалізовані всі 12+ тем практики

✅ Модульна архітектура з 8 окремих файлів

✅ Інтеграція з реальним API (Unsplash)

✅ Три різні режими завантаження даних

✅ Адаптивний дизайн для всіх пристроїв

✅ Система збереження улюблених фото

✅ Професійна обробка помилок

✅ Оптимізація (debouncing, lazy loading, delegation)

✅ Готовність до deployment"""

    tf.text = text
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = GREEN
        paragraph.font.bold = True
        paragraph.space_after = Pt(12)

    # === СЛАЙД 16: Інструкція запуску ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    tf.text = "🚀 Як запустити проєкт"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(7), Inches(4.5))
    tf = content_box.text_frame
    text = """1️⃣ Клонування репозиторію:
   git clone https://github.com/AlexBram003/Practika
   cd Practika

2️⃣ Встановлення залежностей:
   npm install

3️⃣ Запуск dev-сервера:
   npm run dev

4️⃣ Відкрити браузер:
   http://localhost:5173

💡 Спробуйте: пошук "nature", клік на категорію,
   додавання до улюблених, зміна режиму завантаження"""

    tf.text = text
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = DARK_GRAY
        paragraph.space_after = Pt(12)

    # === СЛАЙД 17: Дякую ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = BLUE
    background.line.fill.background()

    thanks_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    tf = thanks_box.text_frame
    tf.text = "Дякую за увагу!"
    p = tf.paragraphs[0]
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    contact_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1.5))
    tf = contact_box.text_frame
    tf.text = "📧 Питання?\n\n📸 Photo Gallery - Брам Олександр, ОП-221"
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = WHITE
        paragraph.alignment = PP_ALIGN.CENTER

    return prs

if __name__ == "__main__":
    print("Створення презентації з візуальними слайдами...")
    prs = create_presentation()
    prs.save('Photo_Gallery_Презентація.pptx')
    print("✅ Презентацію збережено: Photo_Gallery_Презентація.pptx")
    print("\n📝 ВАЖЛИВО: Додайте скріншоти сайту у слайди 2-11!")
    print("   Відкрийте презентацію в PowerPoint та замініть сірі блоки")
    print("   на реальні скріншоти застосунку.")
